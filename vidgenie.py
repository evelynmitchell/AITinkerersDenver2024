import os
import json
from typing import List, Union, Optional

import dspy
import pandas as pd
import pm4py
from pydantic import BaseModel
from twelvelabs import TwelveLabs
from groq import Groq
import chromadb
import requests
from pprint import pprint
import logging
import time
from twelvelabs.exceptions import InternalServerError

# Set up API keys and environment variables
os.environ['API_URL'] = 'https://api.twelvelabs.io/v1.2'
os.environ['TWELVE_LABS_API_KEY'] = 'tlk_2N42Q7B0R5JHV029BKVH51WXV26H'  # It's better to set this outside the script
os.environ['GROQ_API_KEY'] = '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'  # It's better to set this outside the script
TWELVE_LABS_API_KEY = os.getenv("TWELVE_LABS_API_KEY")
GROQ_API_KEY = os.getenv("GROQ_API_KEY")

# Check if API keys are set
if not TWELVE_LABS_API_KEY:
    raise ValueError("TWELVE_LABS_API_KEY is not set or is empty. Please set this environment variable.")
if not GROQ_API_KEY:
    raise ValueError("GROQ_API_KEY is not set or is empty. Please set this environment variable.")

# Initialize clients
twelve_labs_client = TwelveLabs(api_key=TWELVE_LABS_API_KEY)
groq_client = Groq(api_key=GROQ_API_KEY)
chroma_client = chromadb.Client()

class Slide(BaseModel):
    title: str
    content: Union[str, List[str]]
    image_path: Optional[str] = None

class Presentation(BaseModel):
    title: str
    subtitle: str
    slides: List[Slide]
    file_path: Optional[str] = None

    def insert_slide_at(self, slide: Slide, position: int) -> None:
        self.slides.insert(position, slide)

    def to_ppt(self, file_path: str = None):
        if file_path:
            self.file_path = file_path

        from pptx import Presentation as PPTPresentation
        from pptx.util import Inches
        from PIL import Image

        prs = PPTPresentation()

        # Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = self.title
        subtitle.text = self.subtitle

        # Add Slides
        for slide_data in self.slides:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = slide_data.title

            if isinstance(slide_data.content, list):
                content.text = "\n".join(slide_data.content)
            else:
                content.text = slide_data.content

            if slide_data.image_path:
                img = Image.open(slide_data.image_path)
                width, height = img.size
                aspect_ratio = width / height

                max_width = Inches(6)
                max_height = Inches(4.5)

                if aspect_ratio > 1:
                    if width > max_width:
                        width = max_width
                        height = width / aspect_ratio
                else:
                    if height > max_height:
                        height = max_height
                        width = height * aspect_ratio

                left = (prs.slide_width - width) / 2
                top = (prs.slide_height - height) / 2

                slide.shapes.add_picture(slide_data.image_path, left, top, width, height)

        prs.save(file_path)
        print(f"Presentation saved as '{file_path}'")

class VideoAnalysis(BaseModel):
    summary: str
    topics: List[str]
    chapters: List[dict]
    transcript: str

class VidGenieContent(BaseModel):
    title: str
    content_type: str
    sections: List[dict]

class GenerateGlossarySlide(dspy.Signature):
    """
    Generates a glossary slide for the provided presentation JSON string.
    The format should be:

    - {TERM}: {DEFINITION}
    """
    presentation_json_str = dspy.InputField(desc="Presentation JSON string.")
    glossary_slide = dspy.OutputField(desc="Generated glossary slide content. NO FORMATING OR MARKDOWN. PLAINTEXT ONLY", prefix="```plaintext\nGlossary Slide:\n\n")

class GlossaryModule(dspy.Module):
    def __init__(self, **forward_args):
        super().__init__()
        self.forward_args = forward_args
        self.output = None

    def forward(self, presentation_json_str):
        pred = dspy.Predict(GenerateGlossarySlide)
        self.output = pred(presentation_json_str=presentation_json_str).glossary_slide
        return self.output

class VidGenieGenerator:
    def __init__(self):
        self.twelve_labs_client = twelve_labs_client  # Use the global client
        self.index_id = "66aa6336c0cd130695ae109c"  # Make sure this is the correct index ID
        self.api_key = TWELVE_LABS_API_KEY

    def list_indexes(self):
        try:
            indexes = self.twelve_labs_client.index.list()
            print(f"Type of indexes: {type(indexes)}")
            print(f"Structure of indexes: {indexes}")
            
            for index_obj in indexes.root:
                print(f"Index ID: {index_obj.id}, Name: {index_obj.name}, Created At: {index_obj.created_at}, Updated At: {index_obj.updated_at}")
                for engine in index_obj.engines.root:
                    print(f"Engine Name: {engine.name}, Options: {engine.options}")
        except Exception as e:
            print(f"Error retrieving indexes: {e}")

    def analyze_video(self, video_file: str, is_youtube_url: bool = False, max_retries: int = 3, retry_delay: int = 5) -> Optional[VideoAnalysis]:
        logging.basicConfig(level=logging.INFO)
        logger = logging.getLogger(__name__)
        logger.info(f"Analyzing video: {video_file}, is_youtube_url: {is_youtube_url}")
        
        headers = {
            "Accept": "application/json",
            "x-api-key": self.twelve_labs_client.api_key
        }
        
        for attempt in range(max_retries):
            try:
                if is_youtube_url:
                    logger.info("Processing YouTube URL")
                    payload = {
                        "index_id": self.index_id,
                        "url": video_file
                    }
                    response = requests.post("https://api.twelvelabs.io/v1.2/tasks", json=payload, headers=headers)
                else:
                    logger.info("Processing local file")
                    if not os.path.exists(video_file):
                        logger.error(f"Error: The file '{video_file}' does not exist.")
                        return None
                    file_size = os.path.getsize(video_file)
                    logger.info(f"File size: {file_size} bytes")
                    if file_size > 2000000000:  # 2GB limit
                        logger.error("Error: File size exceeds 2GB limit.")
                        return None
                    
                    with open(video_file, 'rb') as file:
                        logger.info(f"File opened successfully: {video_file}")
                        files = {'file': (os.path.basename(video_file), file)}
                        data = {'index_id': self.index_id}
                        response = requests.post("https://api.twelvelabs.io/v1.2/tasks", files=files, data=data, headers=headers)
                
                response.raise_for_status()
                task = response.json()
                logger.info(f"Task created successfully: {task['id']}")
                break  # If successful, break out of the retry loop
            except requests.exceptions.RequestException as e:
                logger.error(f"Attempt {attempt + 1} failed: {e}")
                logger.error(f"Response status code: {response.status_code}")
                logger.error(f"Response headers: {response.headers}")
                logger.error(f"Response content: {response.text}")
                try:
                    error_details = response.json()
                    logger.error(f"Error details: {json.dumps(error_details, indent=2)}")
                except json.JSONDecodeError:
                    logger.error("Could not parse error response as JSON")
                if attempt < max_retries - 1:
                    logger.info(f"Retrying in {retry_delay} seconds...")
                    time.sleep(retry_delay)
                else:
                    logger.error("Max retries reached. Unable to process the video.")
                    return None

        try:
            # Poll for task completion
            while True:
                response = requests.get(f"https://api.twelvelabs.io/v1.2/tasks/{task['id']}", headers=headers)
                response.raise_for_status()
                task_status = response.json()
                if task_status['status'] == 'ready':
                    break
                time.sleep(2)
            
            logger.info(f"Embedding done: {task_status['status']}")
            
            # Retrieve video information
            response = requests.get(f"https://api.twelvelabs.io/v1.2/indexes/{self.index_id}/videos/{task_status['video_id']}", headers=headers)
            response.raise_for_status()
            video_info = response.json()
            
            # Generate summary
            summary_response = requests.post(f"https://api.twelvelabs.io/v1.2/summarize", 
                                             json={"video_id": task_status['video_id'], "type": "summary"}, 
                                             headers=headers)
            summary_response.raise_for_status()
            summary = summary_response.json()['summary']
            
            # Generate gist (topics and chapters)
            gist_response = requests.post(f"https://api.twelvelabs.io/v1.2/generate", 
                                          json={"video_id": task_status['video_id'], "types": ["topic", "chapter"]}, 
                                          headers=headers)
            gist_response.raise_for_status()
            gist = gist_response.json()
            
            # Generate transcript
            transcript_response = requests.post(f"https://api.twelvelabs.io/v1.2/generate", 
                                                json={"video_id": task_status['video_id'], "prompt": "Provide a full transcript of the video"}, 
                                                headers=headers)
            transcript_response.raise_for_status()
            transcript = transcript_response.json()['data']
            
            return VideoAnalysis(
                summary=summary,
                topics=gist['topics'],
                chapters=gist['chapters'],
                transcript=transcript
            )
        except requests.exceptions.RequestException as e:
            logger.error(f"An error occurred while processing the task: {e}")
            import traceback
            traceback.print_exc()
            return None

    def generate_content(self, video_analysis: VideoAnalysis, content_type: str) -> VidGenieContent:
        prompt = f"""
        Based on the following video analysis, generate a {content_type} for an educational course:
        
        Summary: {video_analysis.summary}
        
        Topics: {', '.join(video_analysis.topics)}
        
        Chapters: {video_analysis.chapters}
        
        Transcript: {video_analysis.transcript[:1000]}...  # Truncated for brevity
        
        Please format the {content_type} as a JSON structure with the following format:
        {{
            "title": "Overall title for the {content_type}",
            "sections": [
                {{
                    "title": "Section title",
                    "content": "Section content or bullet points"
                }},
                ...
            ]
        }}
        """
        
        response = groq_client.chat.completions.create(
            messages=[{"role": "user", "content": prompt}],
            model="mixtral-8x7b-32768",
        )
        
        generated_content = response.choices[0].message.content
        parsed_content = json.loads(generated_content)
        
        return VidGenieContent(
            title=parsed_content["title"],
            content_type=content_type,
            sections=parsed_content["sections"]
        )

    def create_syllabus(self, content: VidGenieContent) -> str:
        syllabus = f"# {content.title}\n\n"
        for section in content.sections:
            syllabus += f"## {section['title']}\n{section['content']}\n\n"
        return syllabus

    def create_study_guide(self, content: VidGenieContent) -> str:
        study_guide = f"# {content.title}\n\n"
        for section in content.sections:
            study_guide += f"## {section['title']}\n{section['content']}\n\n"
        return study_guide

    def create_presentation(self, content: VidGenieContent) -> Presentation:
        slides = [Slide(title=section['title'], content=section['content']) for section in content.sections]
        presentation = Presentation(
            title=content.title,
            subtitle="Generated by VidGenie",
            slides=slides
        )
        
        glossary_slide = self.generate_glossary_slide(presentation)
        presentation.insert_slide_at(glossary_slide, position=-1)
        
        return presentation

    def generate_glossary_slide(self, presentation: Presentation) -> Slide:
        presentation_json_str = presentation.json()
        glossary = GlossaryModule()
        glossary_slide_content = glossary.forward(presentation_json_str=presentation_json_str)
        return Slide(title="Glossary of Terms", content=glossary_slide_content)

    def query_video_content(self, query: str) -> List[dict]:
        collection = chroma_client.get_or_create_collection("video_embeddings")
        
        embedding = self.twelve_labs_client.embed.create(
            engine_name="Marengo-retrieval-2.6",
            text=query,
            text_truncate="start",
        )
        
        query_result = collection.query(
            query_embeddings=[embedding.text_embedding.float],
            n_results=2,
        )
        
        return query_result['metadatas'][0]

    def generate_response(self, video_id: str, prompt: str) -> str:
        GENERATE_URL = "https://api.twelvelabs.io/v1.2/generate"
        headers = {"x-api-key": TWELVE_LABS_API_KEY}
        data = {"video_id": video_id, "prompt": prompt}
        
        response = requests.post(GENERATE_URL, headers=headers, json=data)
        return response.json()

def process_video(video_file: str, content_types: List[str], is_youtube_url: bool = False):
    print(f"Processing video: {video_file}, content_types: {content_types}, is_youtube_url: {is_youtube_url}")
    generator = VidGenieGenerator()
    print("Analyzing video")
    video_analysis = generator.analyze_video(video_file, is_youtube_url)
    if video_analysis is None:
        print("Video analysis failed.")
        return None
    
    print("Video analysis completed successfully")
    results = {}

    for content_type in content_types:
        print(f"Generating content for type: {content_type}")
        content = generator.generate_content(video_analysis, content_type)
        if content_type == "syllabus":
            results[content_type] = generator.create_syllabus(content)
        elif content_type == "study_guide":
            results[content_type] = generator.create_study_guide(content)
        elif content_type == "presentation":
            presentation = generator.create_presentation(content)
            results[content_type] = presentation.to_dict()

    print(f"Content generation completed. Results: {results}")
    return results

def main():
    print("Starting main function")
    from dspygen.utils.dspy_tools import init_ol
    init_ol()

    print("Creating VidGenieGenerator instance")
    generator = VidGenieGenerator()
    
    print("Listing indexes")
    generator.list_indexes()

    # Example usage
    video_input = "Building Custom GPTs with GPT Crawler.mp4"
    is_youtube_url = False
    content_types = ["summary", "transcript"]
    
    print(f"Starting video processing with input: {video_input}")
    results = process_video(video_input, content_types, is_youtube_url)
    
    if results is None:
        print("Video processing failed. Exiting.")
        return

    print("Video processing completed successfully")
    print(json.dumps(results, indent=2))

    # Create PowerPoint presentation if requested
    if "presentation" in results:
        print("Creating PowerPoint presentation")
        presentation = Presentation(**results["presentation"])
        presentation.to_ppt(file_path="generated_presentation.pptx")

    # Example of querying video content
    print("Querying video content")
    query_results = generator.query_video_content("custom GPTs")
    print("Query results:", query_results)

    # Example of generating a response based on video content
    if query_results:
        print("Generating response based on video content")
        video_id = query_results[0]['video_id']
        response = generator.generate_response(video_id, "What are the main steps to create a custom GPT?")
        print("Generated response:", response)

    print("Main function completed")

if __name__ == '__main__':
    main()
